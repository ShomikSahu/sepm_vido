import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_vido_presentation(output_path="vido_sepm_presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette - Dark Observatory Technical Theme (Minimalist, High Contrast & SEPM Restraint)
    BG_COLOR = RGBColor(11, 14, 20)       # #0B0E14 Deep Obsidian Background
    CARD_BG = RGBColor(19, 26, 38)        # #131A26 Dark Slate Card
    CARD_BORDER = RGBColor(38, 52, 74)    # #26344A Card Border
    TEXT_MAIN = RGBColor(248, 250, 252)   # #F8FAFC Bright White
    TEXT_MUTED = RGBColor(203, 213, 225)  # #CBD5E1 Slate White
    TEXT_DIM = RGBColor(148, 163, 184)    # #94A3B8 Dim Slate
    ACCENT_CYAN = RGBColor(56, 189, 248)   # #38BDF8 Observatory Cyan (Primary Accent)
    ACCENT_ORANGE = RGBColor(255, 87, 34)  # #FF5722 Secondary Accent
    ACCENT_GOLD = RGBColor(245, 158, 11)   # #F59E0B Highlight Gold
    ACCENT_GREEN = RGBColor(52, 211, 153)  # #34D399 Status Green
    CODE_BG = RGBColor(15, 23, 42)        # #0F172A Code / Diagram Container BG

    blank_layout = prs.slide_layouts[6]

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="SOFTWARE ENGINEERING & PROJECT MANAGEMENT (SEPM)", slide_num=None):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.3))
        tf = cat_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.font.name = "Arial"

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(10.0), Inches(0.6))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MAIN
        p2.font.name = "Arial"

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.fill.background()

        if slide_num is not None:
            num_box = slide.shapes.add_textbox(Inches(11.0), Inches(0.4), Inches(1.533), Inches(0.3))
            tf_n = num_box.text_frame
            tf_n.word_wrap = True
            tf_n.margin_left = tf_n.margin_top = tf_n.margin_right = tf_n.margin_bottom = 0
            pn = tf_n.paragraphs[0]
            pn.text = f"SLIDE {slide_num} / 15"
            pn.alignment = PP_ALIGN.RIGHT
            pn.font.size = Pt(10)
            pn.font.bold = True
            pn.font.color.rgb = TEXT_DIM
            pn.font.name = "Arial"

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # =============================================================
    # SLIDE 1: Title Slide (SEPM Centred)
    # =============================================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.1), Inches(4.5), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = CARD_BG
    badge.line.color.rgb = ACCENT_CYAN
    badge.line.width = Pt(1)
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "SOFTWARE ENGINEERING & PROJECT MANAGEMENT"
    p_b.alignment = PP_ALIGN.CENTER
    p_b.font.size = Pt(10.5)
    p_b.font.bold = True
    p_b.font.color.rgb = ACCENT_CYAN

    tbox = slide1.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.733), Inches(2.2))
    tf1 = tbox.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Volcanic Image/Data Observatory"
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN
    
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "(VIDO)"
    p1_sub.font.size = Pt(34)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = ACCENT_CYAN

    sub_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.0))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_s = tf_sub.paragraphs[0]
    p_s.text = "A Lightweight Scientific Software System for Heterogeneous Terrestrial & Planetary Volcanology Exploration"
    p_s.font.size = Pt(16)
    p_s.font.color.rgb = TEXT_MUTED

    meta_card = add_card(slide1, Inches(0.8), Inches(5.3), Inches(11.733), Inches(1.4))
    tbox_m = slide1.shapes.add_textbox(Inches(1.0), Inches(5.42), Inches(11.333), Inches(1.2))
    tf_m = tbox_m.text_frame
    tf_m.word_wrap = True
    
    pm1 = tf_m.paragraphs[0]
    pm1.text = "Software Engineering & Project Management (SEPM) • Final System Evaluation & Audit"
    pm1.font.size = Pt(14)
    pm1.font.bold = True
    pm1.font.color.rgb = TEXT_MAIN
    
    pm2 = tf_m.add_paragraph()
    pm2.text = "Application Domain: Terrestrial (Earth WGS84) & Planetary Volcanology (Mars, Io, Venus)"
    pm2.font.size = Pt(12)
    pm2.font.color.rgb = ACCENT_CYAN

    pm3 = tf_m.add_paragraph()
    pm3.text = "Implementation Verification: Fully Verified Core MVP • 37 Passing Automated Tests"
    pm3.font.size = Pt(12)
    pm3.font.bold = True
    pm3.font.color.rgb = ACCENT_GREEN

    # =============================================================
    # SLIDE 2: Problem Definition (SE Framing)
    # =============================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "1. Problem Definition: Heterogeneous Scientific Data Management", slide_num=2)

    # Card 1: Core SE Question
    add_card(slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb_p1 = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_p1 = tb_p1.text_frame
    tf_p1.word_wrap = True
    
    p = tf_p1.paragraphs[0]
    p.text = "CORE ENGINEERING QUESTION"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p_q = tf_p1.add_paragraph()
    p_q.space_before = Pt(8)
    p_q.text = "“How can heterogeneous scientific observations be represented, validated, related, and explored in a coherent and maintainable software system?”"
    p_q.font.size = Pt(12)
    p_q.font.italic = True
    p_q.font.color.rgb = TEXT_MAIN

    bullets1 = [
        ("Heterogeneous Scientific Data:", " Multispectral imagery, thermal radiometry, orbital geometry, and point coordinates originating from diverse instruments."),
        ("Application Domain Context:", " Scientific observation datasets across Earth, Mars, Io, and Venus introducing non-standard spatial and temporal constraints."),
        ("Software Engineering Framing:", " Volcanology is the application domain; software design, architecture, data modeling, and verification are the primary subject.")
    ]
    for b_title, b_desc in bullets1:
        p_b = tf_p1.add_paragraph()
        p_b.space_before = Pt(8)
        r1 = p_b.add_run()
        r1.text = "• " + b_title
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = TEXT_MAIN
        r2 = p_b.add_run()
        r2.text = b_desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # Card 2: SE Challenges
    add_card(slide2, Inches(6.8), Inches(1.6), Inches(5.733), Inches(5.2))
    tb_p2 = slide2.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_p2 = tb_p2.text_frame
    tf_p2.word_wrap = True

    p = tf_p2.paragraphs[0]
    p.text = "PRIMARY SOFTWARE ENGINEERING CHALLENGES"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE

    bullets2 = [
        ("Coordinate Frame Disparities:", " Standard Earth WGS84 (-180° to +180°) vs. Planetary IAU (0° to 360° Positive East) conventions."),
        ("Schema Rigidity & Bloat:", " Forcing multi-sensor attributes into flat relational SQL schemas causes sparse, brittle tables with excessive NULL columns."),
        ("Observation / Event Conflation:", " Conflating observational snapshot evidence with physical eruptive events degrades domain modeling integrity."),
        ("System Maintainability:", " Delivering strict domain validation and clean data access without heavy, complex framework dependencies.")
    ]
    for b_title, b_desc in bullets2:
        p_b = tf_p2.add_paragraph()
        p_b.space_before = Pt(10)
        r1 = p_b.add_run()
        r1.text = "• " + b_title
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = TEXT_MAIN
        r2 = p_b.add_run()
        r2.text = b_desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 3: Requirements Engineering
    # =============================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "2. Requirements Engineering: Functional & Non-Functional Specifications", slide_num=3)

    # Left: Functional Requirements Card
    add_card(slide3, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb_fr = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_fr = tb_fr.text_frame
    tf_fr.word_wrap = True

    p = tf_fr.paragraphs[0]
    p.text = "FUNCTIONAL REQUIREMENTS (FR)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    fr_items = [
        ("FR-01 Observation Management:", " Store, update, and manage observation records across celestial bodies."),
        ("FR-02 Multi-Criteria Retrieval:", " Query observations by date range, facet category, sensor source, and system."),
        ("FR-03 Event Association:", " Relate observations to physical eruptive events using explicit relationship tags."),
        ("FR-04 Spatial Exploration:", " Plot spatial markers on Earth Leaflet maps or 0°-360° planetary grids with fallback rules."),
        ("FR-05 Timeline Exploration:", " Render synchronized dual-lane chronological timeline with ongoing event support."),
        ("FR-06 Domain Validation:", " Enforce metadata facet validation and coordinate boundary rules prior to persistence.")
    ]
    for f_title, f_desc in fr_items:
        pf = tf_fr.add_paragraph()
        pf.space_before = Pt(5)
        r1 = pf.add_run()
        r1.text = "• " + f_title
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = TEXT_MAIN
        r2 = pf.add_run()
        r2.text = f_desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # Right: Non-Functional Requirements Card
    add_card(slide3, Inches(6.8), Inches(1.6), Inches(5.733), Inches(5.2))
    tb_nfr = slide3.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_nfr = tb_nfr.text_frame
    tf_nfr.word_wrap = True

    p = tf_nfr.paragraphs[0]
    p.text = "NON-FUNCTIONAL REQUIREMENTS (NFR TARGETS)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE

    nfr_items = [
        ("NFR-01 Maintainability:", " Layered 5-tier architecture enforcing separation of concerns across presentation, API, services, repos, and DB."),
        ("NFR-02 Extensibility:", " Modular facet schemas allowing new observation types without database schema migrations."),
        ("NFR-03 Testability:", " Decoupled domain service layer enabling isolated automated testing (core domain logic test suite)."),
        ("NFR-04 Data Integrity:", " Schema boundary validation (Pydantic) and database relational transaction safety."),
        ("NFR-05 Performance & Modularity:", " Lightweight architecture with zero heavy build tool overhead and low query latencies.")
    ]
    for n_title, n_desc in nfr_items:
        pn = tf_nfr.add_paragraph()
        pn.space_before = Pt(8)
        r1 = pn.add_run()
        r1.text = "• " + n_title
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = TEXT_MAIN
        r2 = pn.add_run()
        r2.text = n_desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 4: Use Case Analysis + UML Use Case Diagram
    # =============================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "3. Use Case Analysis & Mandatory UML Use Case Diagram", slide_num=4)

    # Left: Engineering Explanation & 6 Use Cases List
    add_card(slide4, Inches(0.8), Inches(1.6), Inches(5.4), Inches(5.2))
    tb_uc = slide4.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_uc = tb_uc.text_frame
    tf_uc.word_wrap = True

    p = tf_uc.paragraphs[0]
    p.text = "USE CASE MODEL & USER INTERACTIONS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p_exp = tf_uc.add_paragraph()
    p_exp.space_before = Pt(6)
    p_exp.text = "“The use-case model translates functional requirements into observable interactions between the user and VIDO.”"
    p_exp.font.size = Pt(11.5)
    p_exp.font.italic = True
    p_exp.font.color.rgb = TEXT_MAIN

    uc_list = [
        ("UC-01 Browse Celestial Systems:", " Select target body & explore system profiles."),
        ("UC-02 Search & Filter Observations:", " Query by date, facet category, or source."),
        ("UC-03 Inspect Metadata Payloads:", " View validated composite metadata facets."),
        ("UC-04 Associate Observations with Events:", " Link evidence snapshots to eruptive events."),
        ("UC-05 Explore Synchronized Timeline:", " Inspect dual-lane chronological feed."),
        ("UC-06 Explore Spatial Coordinates:", " Render markers with spatial fallback rules.")
    ]
    for u_title, u_desc in uc_list:
        pu = tf_uc.add_paragraph()
        pu.space_before = Pt(5)
        r1 = pu.add_run()
        r1.text = "• " + u_title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = ACCENT_CYAN
        r2 = pu.add_run()
        r2.text = u_desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # Right: Embedded PlantUML Use Case Diagram Image
    add_card(slide4, Inches(6.4), Inches(1.6), Inches(6.133), Inches(5.2), bg_color=CODE_BG)
    uc_img_path = os.path.join("docs", "plantuml", "use_case_diagram.png")
    if os.path.exists(uc_img_path):
        slide4.shapes.add_picture(uc_img_path, Inches(6.6), Inches(1.75), width=Inches(5.733))

    # =============================================================
    # SLIDE 5: Requirements Traceability Matrix (Verified Codebase Names)
    # =============================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "4. Requirements Traceability: FR → UC → API → Service → Test", slide_num=5)

    # Top Pipeline Header
    add_card(slide5, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.8))
    tb_t1 = slide5.shapes.add_textbox(Inches(1.0), Inches(1.58), Inches(11.333), Inches(0.6))
    tf_t1 = tb_t1.text_frame
    tf_t1.word_wrap = True
    p = tf_t1.paragraphs[0]
    p.text = "END-TO-END SYSTEM TRACEABILITY PIPELINE"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p_flow = tf_t1.add_paragraph()
    p_flow.space_before = Pt(3)
    p_flow.text = "Requirement (FR)   →   Use Case (UC)   →   API Endpoint   →   Service Layer   →   Repository / DB   →   Pytest Verification"
    p_flow.font.size = Pt(11)
    p_flow.font.bold = True
    p_flow.font.color.rgb = TEXT_MAIN

    # Traceability Rows strictly audited against verified codebase classes
    t_rows = [
        ("FR-01 Observation Management", "UC-02 / UC-03", "POST /api/v1/observations", "ValidationService", "ObservationRepository", "test_models.py, test_api.py"),
        ("FR-02 Multi-Criteria Retrieval", "UC-02", "GET /api/v1/observations", "ValidationService", "ObservationRepository", "test_api.py"),
        ("FR-03 Event Association", "UC-04", "POST /api/v1/observation-event-links", "EventService", "ObservationEventLinkRepository", "test_domain_rules.py, test_services.py"),
        ("FR-04 Spatial Exploration", "UC-06", "GET /api/v1/systems/{id}/spatial", "SpatialService / CoordinateService", "VolcanicSystemRepository / ObservationRepo", "test_domain_rules.py, test_repositories.py"),
        ("FR-05 Timeline Exploration", "UC-05", "GET /api/v1/systems/{id}/timeline", "TimelineService", "VolcanicEventRepository / ObservationRepo", "test_services.py, test_api.py"),
        ("FR-06 Domain Validation", "UC-02 / UC-03", "Pydantic API Schemas", "ValidationService / CoordinateService", "Domain Model Constraints", "test_models.py, test_domain_rules.py")
    ]

    card_h5 = Inches(0.72)
    for idx, (fr, uc, api, srv, repo, test) in enumerate(t_rows):
        top_pos = Inches(2.45) + Inches(idx * 0.78)
        add_card(slide5, Inches(0.8), top_pos, Inches(11.733), card_h5)
        
        tb_r = slide5.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.08), Inches(11.333), card_h5 - Inches(0.16))
        tf_r = tb_r.text_frame
        tf_r.word_wrap = True
        
        pr1 = tf_r.paragraphs[0]
        pr1.text = f"{fr}  |  Use Case: {uc}  |  API: {api}"
        pr1.font.size = Pt(11.5)
        pr1.font.bold = True
        pr1.font.color.rgb = ACCENT_CYAN if idx % 2 == 0 else ACCENT_ORANGE
        
        pr2 = tf_r.add_paragraph()
        pr2.space_before = Pt(2)
        pr2.text = f"Service: {srv}   •   Repo: {repo}   •   Verification: {test}"
        pr2.font.size = Pt(10.5)
        pr2.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 6: Software Design + UML Class Diagram
    # =============================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "5. Software Design & Mandatory UML Class Diagram", slide_num=6)

    # Left Column: Domain Entity Roles & Concise Engineering Explanation
    add_card(slide6, Inches(0.8), Inches(1.6), Inches(5.4), Inches(5.2))
    tb_sd = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_sd = tb_sd.text_frame
    tf_sd.word_wrap = True

    p = tf_sd.paragraphs[0]
    p.text = "DOMAIN CLASS MODEL & RELATIONSHIPS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p_exp6 = tf_sd.add_paragraph()
    p_exp6.space_before = Pt(6)
    p_exp6.text = "“Separating observations from physical events preserves distinct lifecycles and attributes while ObservationEventLink provides explicit association.”"
    p_exp6.font.size = Pt(11)
    p_exp6.font.italic = True
    p_exp6.font.color.rgb = TEXT_MAIN

    class_roles = [
        ("CelestialBody (1 → * VolcanicSystem):", " Planet/Moon container establishing coordinate convention."),
        ("VolcanicSystem (1 → * Observation):", " Volcanic feature defining geographic base coordinates."),
        ("ObservationSource (1 → * Observation):", " Sensor platform or operating scientific agency."),
        ("Observation (1 → * Link):", " Point-in-time observational evidence snapshot with metadata."),
        ("VolcanicEvent (1 → * Link):", " Physical eruptive episode with duration (start_time, end_time)."),
        ("ObservationEventLink:", " Explicit junction entity connecting observations to events.")
    ]
    for c_title, c_desc in class_roles:
        pc = tf_sd.add_paragraph()
        pc.space_before = Pt(4)
        r1 = pc.add_run()
        r1.text = "• " + c_title + " "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = ACCENT_CYAN
        r2 = pc.add_run()
        r2.text = c_desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_MUTED

    # Right Column: Embedded PlantUML Class Diagram Image
    add_card(slide6, Inches(6.4), Inches(1.6), Inches(6.133), Inches(5.2), bg_color=CODE_BG)
    class_img_path = os.path.join("docs", "plantuml", "class_diagram.png")
    if os.path.exists(class_img_path):
        slide6.shapes.add_picture(class_img_path, Inches(7.5), Inches(1.7), height=Inches(5.0))

    # =============================================================
    # SLIDE 7: Composite Metadata / Facet Design
    # =============================================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_bg(slide7)
    add_header(slide7, "6. Composite Metadata Facet Architecture", slide_num=7)

    # Left Column: Concept & Facet Hierarchy
    add_card(slide7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb_c1 = slide7.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True

    p = tf_c1.paragraphs[0]
    p.text = "MODULAR FACET COMPOSITION PATTERN"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    facets = [
        ("Observation Facet Structure:", " Observation entity encapsulates non-mutually-exclusive metadata facets rather than rigid subclassing."),
        ("IMAGE Facet:", " Spectral band, spatial resolution (m), cloud cover %, file format, sun elevation angle."),
        ("THERMAL Facet:", " Brightness temperature (K), ambient temp (K), thermal flux (MW), anomaly flag."),
        ("PLANETARY_ORBITAL Facet:", " Spacecraft altitude (km), solar incidence angle, emission angle, target planetary datum."),
        ("Validation Guarantee:", " Pydantic v2 schemas reject malformed or out-of-bounds payloads at API boundaries.")
    ]
    for f_title, f_desc in facets:
        pf = tf_c1.add_paragraph()
        pf.space_before = Pt(6)
        r1 = pf.add_run()
        r1.text = "• " + f_title + " "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = TEXT_MAIN
        r2 = pf.add_run()
        r2.text = f_desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # Right Column: Actual Olympus Mons Composite Payload
    add_card(slide7, Inches(6.8), Inches(1.6), Inches(5.733), Inches(5.2), bg_color=CODE_BG)
    tb_c2 = slide7.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True

    p = tf_c2.paragraphs[0]
    p.text = "REPRESENTATIVE COMPOSITE JSON PAYLOAD"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD

    json_snippet = """{
  "active_facets": [
    "IMAGE", "THERMAL", "PLANETARY_ORBITAL"
  ],
  "image_metadata": {
    "spectral_band": "NEAR_INFRARED",
    "spatial_resolution_m": 10.0,
    "cloud_cover_percentage": 0.0
  },
  "thermal_metadata": {
    "brightness_temperature_kelvin": 245.5,
    "thermal_flux_mw": 12.4,
    "anomaly_flag": false
  },
  "orbital_metadata": {
    "spacecraft_altitude_km": 310.5,
    "solar_incidence_angle_deg": 62.1,
    "target_planetary_datum": "IAU_MARS_2000"
  }
}"""
    p_code = tf_c2.add_paragraph()
    p_code.space_before = Pt(6)
    p_code.text = json_snippet
    p_code.font.size = Pt(9.5)
    p_code.font.name = "Consolas"
    p_code.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 8: System Architecture (Corrected Service Layer Role)
    # =============================================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_bg(slide8)
    add_header(slide8, "7. System Architecture: 5-Layer Software Pipeline", slide_num=8)

    layers = [
        ("Presentation Layer", "Vanilla HTML5 / CSS3 / ES6 JS", "Renders responsive web UI, Leaflet Earth maps, Canvas 2D planetary grids, and interactive timeline visualizers without heavy frontend build tool overhead.", ACCENT_CYAN),
        ("API Layer", "Python FastAPI REST Server", "Exposes REST endpoints (/api/v1), serializes request/response payloads, executes Pydantic schema boundary validation, handles HTTP status routing, and serves Swagger OpenAPI docs.", ACCENT_ORANGE),
        ("Service Layer — Business & Domain Logic", "Domain & Business Services", "Executes domain/business validation, coordinate resolution, spatial fallback rules, event rules, and timeline aggregation.", ACCENT_GOLD),
        ("Repository Layer", "Data Access Abstraction", "Encapsulates database SQL operations, providing clean, decoupled interfaces for systems, observations, events, and association links.", ACCENT_CYAN),
        ("Database Layer", "SQLite 3 Relational Engine", "Provides persistent, ACID-compliant storage with SQLite 3 JSON1 support for composite observation metadata facets.", ACCENT_GREEN)
    ]

    top_pos8 = Inches(1.6)
    card_h8 = Inches(0.95)

    for idx, (l_title, l_tech, l_desc, l_color) in enumerate(layers):
        c_top = top_pos8 + Inches(idx * 1.05)
        add_card(slide8, Inches(0.8), c_top, Inches(11.733), card_h8)
        
        tb_l = slide8.shapes.add_textbox(Inches(1.0), c_top + Inches(0.12), Inches(3.7), Inches(0.7))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        
        p1 = tf_l.paragraphs[0]
        p1.text = l_title
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = l_color
        
        p2 = tf_l.add_paragraph()
        p2.text = l_tech
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_DIM

        div = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), c_top + Inches(0.15), Inches(0.02), Inches(0.65))
        div.fill.solid()
        div.fill.fore_color.rgb = CARD_BORDER
        div.line.fill.background()

        tb_r = slide8.shapes.add_textbox(Inches(5.0), c_top + Inches(0.15), Inches(7.3), Inches(0.65))
        tf_r = tb_r.text_frame
        tf_r.word_wrap = True
        p3 = tf_r.paragraphs[0]
        p3.text = l_desc
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = TEXT_MAIN

    # =============================================================
    # SLIDE 9: Design Decisions & Trade-offs
    # =============================================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_bg(slide9)
    add_header(slide9, "8. Design Decisions & Engineering Trade-offs", slide_num=9)

    decisions = [
        ("Relational Core + JSON Metadata", "Combines SQL relational structure with SQLite JSON1 metadata fields. Prevents schema rigidity while retaining query capabilities.", ACCENT_CYAN),
        ("Decoupled Obs & Event Entities", "Models observations (evidence snapshots) separately from events (phenomena). Prevents data loss when linking multi-event observations.", ACCENT_ORANGE),
        ("Body-Specific Coordinate Handling", "Enforces Earth WGS84 (-180°..+180°) vs. Planetary IAU (0°..360° Positive East) conventions. Prevents spatial projection errors.", ACCENT_GOLD),
        ("Explicit Spatial Fallback Rules", "Applies a 3-step fallback pipeline (Obs Coords → Volcano Coords → UNLOCATED). Prevents coordinate fabrication.", ACCENT_CYAN),
        ("Service / Repository Abstraction", "Decouples business logic from persistence logic. Prevents SQL leaks into API controllers and enables unit test isolation.", ACCENT_GREEN),
        ("Conservative MVP Scope Boundary", "Focuses on robust core domain modeling, API, and validation rather than unverified complex modules. Prevents scope creep.", ACCENT_ORANGE)
    ]

    col_w9 = Inches(3.64)
    row_h9 = Inches(2.4)
    offsets9 = [
        (Inches(0.8), Inches(1.6)), (Inches(4.84), Inches(1.6)), (Inches(8.88), Inches(1.6)),
        (Inches(0.8), Inches(4.4)), (Inches(4.84), Inches(4.4)), (Inches(8.88), Inches(4.4))
    ]

    for idx, (title, desc, color) in enumerate(decisions):
        left, top = offsets9[idx]
        add_card(slide9, left, top, col_w9, row_h9)
        tb = slide9.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), col_w9 - Inches(0.4), row_h9 - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = color
        
        pd = tf.add_paragraph()
        pd.space_before = Pt(6)
        pd.text = desc
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 10: Validation & Business Rules
    # =============================================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_bg(slide10)
    add_header(slide10, "9. Validation & Business Rules Architecture", slide_num=10)

    # Card 1: API / Pydantic Validation
    add_card(slide10, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb_v1 = slide10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_v1 = tb_v1.text_frame
    tf_v1.word_wrap = True

    p = tf_v1.paragraphs[0]
    p.text = "TIER 1: API SCHEMA VALIDATION (PYDANTIC)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    v1_points = [
        ("Boundary Schema Enforcement:", " Validates request JSON payloads at API entry points before database invocation."),
        ("Attribute Type Checking:", " Guarantees correct string, float, integer, and datetime field formatting."),
        ("Facet Payload Validation:", " Validates nested image_metadata, thermal_metadata, and orbital_metadata objects."),
        ("Value Range Rules:", " Rejects negative Kelvin temperatures, out-of-range sensor angles, or invalid file extensions.")
    ]
    for title, desc in v1_points:
        pv = tf_v1.add_paragraph()
        pv.space_before = Pt(8)
        r1 = pv.add_run()
        r1.text = "• " + title + " "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = TEXT_MAIN
        r2 = pv.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # Card 2: Domain / Business Logic Validation
    add_card(slide10, Inches(6.8), Inches(1.6), Inches(5.733), Inches(5.2))
    tb_v2 = slide10.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_v2 = tb_v2.text_frame
    tf_v2.word_wrap = True

    p = tf_v2.paragraphs[0]
    p.text = "TIER 2: DOMAIN & BUSINESS RULE VALIDATION"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE

    v2_points = [
        ("Body Coordinate Boundary Checks:", " Enforces Earth WGS84 (-180°..+180°) vs. Planetary IAU (0°..360° Positive East) bounds."),
        ("Temporal Event Sequence Rules:", " Ensures event end_time cannot precede start_time (when end_time is provided)."),
        ("Relationship Classification Rules:", " Enforces valid Observation ↔ Event tags (PRE_ERUPTIVE, CO_ERUPTIVE, POST_ERUPTIVE, UNRELATED)."),
        ("Spatial Fallback Resolution:", " Classifies coordinates into OBSERVATION, VOLCANO_FALLBACK, or UNLOCATED states cleanly.")
    ]
    for title, desc in v2_points:
        pv = tf_v2.add_paragraph()
        pv.space_before = Pt(8)
        r1 = pv.add_run()
        r1.text = "→ " + title + " "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = TEXT_MAIN
        r2 = pv.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 11: Development Process (Simplified SDLC Phases)
    # =============================================================
    slide11 = prs.slides.add_slide(blank_layout)
    add_bg(slide11)
    add_header(slide11, "10. Development Process: Phased SDLC Workflow", slide_num=11)

    sdlc_phases = [
        ("1. Problem Definition & SRS", "Formalized engineering problem statement and compiled Software Requirements Specification.", ACCENT_CYAN),
        ("2. Requirements Audit", "Audited domain parameters and mapped functional vs. non-functional requirement targets.", ACCENT_ORANGE),
        ("3. Data Layer & Entities", "Designed SQLite relational tables, JSON metadata fields, and Pydantic validation schemas.", ACCENT_GOLD),
        ("4. Business Logic Services", "Implemented coordinate resolution, spatial fallback rules, event rules, and timeline aggregation.", ACCENT_CYAN),
        ("5. API Layer", "Exposed REST endpoints (/api/v1) with OpenAPI Swagger docs and robust HTTP error routing.", ACCENT_GREEN),
        ("6. Frontend Interface", "Built responsive web UI with Leaflet maps, Canvas 2D planetary grid, and timeline visualizers.", ACCENT_ORANGE),
        ("7. Automated Verification", "Developed pytest test suite covering API routes, domain rules, schemas, and repositories.", ACCENT_CYAN),
        ("8. Presentation & QC", "Audited implementation claims against codebase evidence and prepared evaluation materials.", ACCENT_GREEN)
    ]

    col_w11 = Inches(2.7)
    row_h11 = Inches(2.4)
    offsets11 = [
        (Inches(0.8), Inches(1.6)), (Inches(3.8), Inches(1.6)), (Inches(6.8), Inches(1.6)), (Inches(9.8), Inches(1.6)),
        (Inches(0.8), Inches(4.4)), (Inches(3.8), Inches(4.4)), (Inches(6.8), Inches(4.4)), (Inches(9.8), Inches(4.4))
    ]

    for idx, (title, desc, color) in enumerate(sdlc_phases):
        left, top = offsets11[idx]
        add_card(slide11, left, top, col_w11, row_h11)
        tb = slide11.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), col_w11 - Inches(0.3), row_h11 - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(11.5)
        pt.font.bold = True
        pt.font.color.rgb = color
        
        pd = tf.add_paragraph()
        pd.space_before = Pt(6)
        pd.text = desc
        pd.font.size = Pt(10)
        pd.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 12: Scope & Risk Management
    # =============================================================
    slide12 = prs.slides.add_slide(blank_layout)
    add_bg(slide12)
    add_header(slide12, "11. Scope Control & Risk Management Strategy", slide_num=12)

    # Card 1: Deliberate Scope Control
    add_card(slide12, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb_sc1 = slide12.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_sc1 = tb_sc1.text_frame
    tf_sc1.word_wrap = True

    p = tf_sc1.paragraphs[0]
    p.text = "DELIBERATE SCOPE CONTROL"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    scope_points = [
        ("Implemented Core MVP Scope:", " Unified domain model, multi-body CRS, composite facet validation, 5-layer pipeline, dual-lane timeline, spatial fallback visualizer, REST API, 37 automated tests."),
        ("Explicit Non-Implemented Future Scope:", " Streaming satellite telemetry pipelines, computer vision thermal anomaly workers, GeoServer WMS integration, user authentication/authorization."),
        ("Scope Boundary Principle:", " Maintaining strict boundaries prevented unfinished features from degrading core software quality.")
    ]
    for title, desc in scope_points:
        ps = tf_sc1.add_paragraph()
        ps.space_before = Pt(8)
        r1 = ps.add_run()
        r1.text = "• " + title + " "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = TEXT_MAIN
        r2 = ps.add_run()
        r2.text = desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # Card 2: Risk Identification & Mitigation
    add_card(slide12, Inches(6.8), Inches(1.6), Inches(5.733), Inches(5.2))
    tb_sc2 = slide12.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_sc2 = tb_sc2.text_frame
    tf_sc2.word_wrap = True

    p = tf_sc2.paragraphs[0]
    p.text = "RISK MITIGATION MATRIX"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ORANGE

    risks = [
        ("Heterogeneous Data Distortion Risk:", " Mitigated by Composite Metadata Facet validation pattern."),
        ("Coordinate Frame Mismatch Risk:", " Mitigated by CelestialBody explicit CRS conventions (WGS84 vs. IAU)."),
        ("Observation/Event Ambiguity Risk:", " Mitigated by explicit ObservationEventLink junction modeling."),
        ("Scope Creep Risk:", " Mitigated by strict SRS specification and phase verification."),
        ("Regression & Code Fragility Risk:", " Mitigated by 37 passing automated pytest unit/integration tests.")
    ]
    for r_title, r_desc in risks:
        pr = tf_sc2.add_paragraph()
        pr.space_before = Pt(8)
        r1 = pr.add_run()
        r1.text = "⚙ " + r_title + " "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = TEXT_MAIN
        r2 = pr.add_run()
        r2.text = r_desc
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 13: Testing & Verification
    # =============================================================
    slide13 = prs.slides.add_slide(blank_layout)
    add_bg(slide13)
    add_header(slide13, "12. Testing & System Verification Results", slide_num=13)

    # Banner: 37/37 Automated Tests Passed
    banner = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = CODE_BG
    banner.line.color.rgb = ACCENT_GREEN
    banner.line.width = Pt(1.5)
    tf_bn = banner.text_frame
    p_bn = tf_bn.paragraphs[0]
    p_bn.text = "✓ 37 / 37 AUTOMATED TESTS PASSED"
    p_bn.alignment = PP_ALIGN.CENTER
    p_bn.font.size = Pt(18)
    p_bn.font.bold = True
    p_bn.font.color.rgb = ACCENT_GREEN

    p_sub_bn = tf_bn.add_paragraph()
    p_sub_bn.text = "Fully verified core MVP implementation across API endpoints, domain rules, metadata schemas, repositories, and services."
    p_sub_bn.alignment = PP_ALIGN.CENTER
    p_sub_bn.font.size = Pt(11)
    p_sub_bn.font.color.rgb = TEXT_MUTED

    test_areas = [
        ("API Endpoint Integration (15 tests)", "API integration tests covering catalog, observation, event, spatial, and timeline REST endpoints.", ACCENT_CYAN),
        ("Domain Rules & Boundaries (5 tests)", "Validates coordinate boundary enforcement (Earth vs. Mars/Io) and temporal sequence validity.", ACCENT_ORANGE),
        ("Model & Facet Schemas (6 tests)", "Tests Pydantic metadata facet validation (IMAGE, THERMAL, PLANETARY_ORBITAL) rejecting illegal payloads.", ACCENT_GOLD),
        ("Repository & Persistence (5 tests)", "Verifies SQLite 3 SQL/JSON data access, insertion transactions, and complex query filtering.", ACCENT_CYAN),
        ("Business Service Logic (5 tests)", "Verifies timeline aggregation logic, ongoing event handling (NULL end_time), and spatial fallback resolution.", ACCENT_GREEN),
        ("Seed Data Verification (1 test)", "Ensures initial seed dataset across Earth, Mars, Io, and Venus populates without relational errors.", ACCENT_ORANGE)
    ]

    col_w13 = Inches(3.64)
    row_h13 = Inches(2.1)
    offsets13 = [
        (Inches(0.8), Inches(2.6)), (Inches(4.84), Inches(2.6)), (Inches(8.88), Inches(2.6)),
        (Inches(0.8), Inches(4.9)), (Inches(4.84), Inches(4.9)), (Inches(8.88), Inches(4.9))
    ]

    for idx, (t_title, t_desc, color) in enumerate(test_areas):
        left, top = offsets13[idx]
        add_card(slide13, left, top, col_w13, row_h13)
        tb = slide13.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), col_w13 - Inches(0.4), row_h13 - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        pt = tf.paragraphs[0]
        pt.text = t_title
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = color
        
        pd = tf.add_paragraph()
        pd.space_before = Pt(6)
        pd.text = t_desc
        pd.font.size = Pt(10.5)
        pd.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 14: Implementation Demonstration
    # =============================================================
    slide14 = prs.slides.add_slide(blank_layout)
    add_bg(slide14)
    add_header(slide14, "13. Implementation Demonstration: Software Evidence", slide_num=14)

    screenshots = [
        ("Mount Etna / Earth", "mount_etna_earth.png", "Leaflet Earth Map view rendering geographic spatial representation & eruptive timeline feed."),
        ("Loki Patera / Io", "loki_patera_io.png", "Canvas 2D Planetary Grid plotting Io thermal observations in 0°-360° Positive East CRS."),
        ("Olympus Mons / Mars", "olympus_mons_mars.png", "Composite Metadata Inspector modal rendering active IMAGE, THERMAL & PLANETARY_ORBITAL facets.")
    ]

    sw = Inches(3.64)
    sh = Inches(5.2)
    s_lefts = [Inches(0.8), Inches(4.84), Inches(8.88)]

    for idx, (title, img_name, desc) in enumerate(screenshots):
        left = s_lefts[idx]
        add_card(slide14, left, Inches(1.6), sw, sh)
        
        tbox = slide14.shapes.add_textbox(left + Inches(0.15), Inches(1.72), sw - Inches(0.3), Inches(0.35))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN if idx != 2 else ACCENT_GOLD
        p.alignment = PP_ALIGN.CENTER

        img_path = os.path.join("docs", "screenshots", img_name)
        if os.path.exists(img_path):
            img_shape = slide14.shapes.add_picture(img_path, left + Inches(0.15), Inches(2.15), width=sw - Inches(0.3))
            
            p_border = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.14), Inches(2.14), sw - Inches(0.28), img_shape.height + Inches(0.02))
            p_border.fill.background()
            p_border.line.color.rgb = CARD_BORDER
            p_border.line.width = Pt(1)

            cap_top = Inches(2.25) + img_shape.height
            cbox = slide14.shapes.add_textbox(left + Inches(0.2), cap_top, sw - Inches(0.4), Inches(1.2))
            tf_c = cbox.text_frame
            tf_c.word_wrap = True
            tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
            pc = tf_c.paragraphs[0]
            pc.text = desc
            pc.font.size = Pt(10.5)
            pc.font.color.rgb = TEXT_MUTED

    # =============================================================
    # SLIDE 15: Team, Limitations & Conclusion (Corrected Restrained Phrasing)
    # =============================================================
    slide15 = prs.slides.add_slide(blank_layout)
    add_bg(slide15)
    add_header(slide15, "14. Team Contributions, Limitations & SEPM Conclusion", slide_num=15)

    # Card 1: Team Contributions & Scope Limitations
    add_card(slide15, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tb_tc = slide15.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_tc = tb_tc.text_frame
    tf_tc.word_wrap = True

    p = tf_tc.paragraphs[0]
    p.text = "TEAM CONTRIBUTIONS & LIMITATIONS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    team_members = [
        ("Technical Implementation:", " Architecture, database schema, FastAPI backend, business validation services, frontend UI, and pytest test suite."),
        ("Documentation & SRS:", " Software Requirements Specification (SRS), domain requirements audit, system documentation, and architecture specifications."),
        ("Presentation & QC Audit:", " Presentation deck, visual alignment, speaker notes, and quality audit."),
        ("Current Scope Limitations:", " Manual dataset seeding; non-automated satellite ingestion; local SQLite storage rather than distributed database.")
    ]
    for m_title, m_desc in team_members:
        pm = tf_tc.add_paragraph()
        pm.space_before = Pt(6)
        r1 = pm.add_run()
        r1.text = "• " + m_title + " "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = TEXT_MAIN
        r2 = pm.add_run()
        r2.text = m_desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_MUTED

    # Card 2: SEPM Conclusion Statement
    add_card(slide15, Inches(6.8), Inches(1.6), Inches(5.733), Inches(5.2))
    tb_c15 = slide15.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.8))
    tf_c15 = tb_c15.text_frame
    tf_c15.word_wrap = True

    p = tf_c15.paragraphs[0]
    p.text = "SEPM CONCLUSION"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    p_conc = tf_c15.add_paragraph()
    p_conc.space_before = Pt(8)
    p_conc.text = "“VIDO demonstrates the design and implementation of a modular scientific information system capable of representing heterogeneous observations, enforcing domain rules, relating observations to events, exposing functionality through a layered API, and supporting interactive exploration with automated verification.”"
    p_conc.font.size = Pt(11.5)
    p_conc.font.italic = True
    p_conc.font.color.rgb = TEXT_MAIN

    future_points = [
        ("Layered Architecture:", " Layered pipeline enforces separation of concerns and supports data integrity and testability."),
        ("Automated Verification:", " 37 / 37 automated tests verify core domain rules, APIs, and repositories."),
        ("Future Scope Roadmap:", " Telemetry streaming pipelines, CV anomaly detection workers, GeoServer WMS integration, and user authorization.")
    ]
    for c_title, c_desc in future_points:
        pc = tf_c15.add_paragraph()
        pc.space_before = Pt(8)
        r1 = pc.add_run()
        r1.text = "✓ " + c_title + " "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = ACCENT_GREEN
        r2 = pc.add_run()
        r2.text = c_desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_MUTED

    try:
        prs.save(output_path)
        print(f"Presentation successfully saved to {output_path}")
    except PermissionError:
        alt_path = "vido_sepm_presentation_updated.pptx"
        prs.save(alt_path)
        print(f"Target file {output_path} is currently locked by PowerPoint. Presentation saved to {alt_path}")

if __name__ == "__main__":
    build_vido_presentation()
